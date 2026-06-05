"""保定オリエンテーション資料を PPTX として生成し、PDF へ変換する。

写真ページには layout.py で生成した 16:9 画像をそのまま全面差し込みする。
テキストページは template.DECK の文面から組み立てる。
"""

from __future__ import annotations

import shutil
import subprocess
import tempfile
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .template import DECK
from .slides import render_deck

# 16:9 スライド寸法
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# 配色（layout.py と同じトーン）
BG = RGBColor(0xF6, 0xF3, 0xF0)
ACCENT = RGBColor(0xD8, 0xA8, 0xB2)
TEXT = RGBColor(0x4A, 0x3C, 0x3A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# 日本語フォント名（PDF変換に使う LibreOffice 環境に存在するものを既定に）
JP_FONT = "IPAGothic"


def _blank_slide(prs: Presentation):
    """背景を塗った白紙スライドを追加して返す。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        0, 0, prs.slide_width, prs.slide_height,
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    # 背景は最背面へ
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return slide


def _textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def _set_run(run, text, size, color, bold=False, font=JP_FONT):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = color


def _add_title_slide(prs, spec):
    slide = _blank_slide(prs)
    tf = _textbox(slide, Inches(1), Inches(2.2), Inches(11.3), Inches(3),
                  anchor=MSO_ANCHOR.MIDDLE)
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    _set_run(p0.add_run(), spec["eyebrow"], 22, ACCENT)
    for i, line in enumerate(spec["title"].split("\n")):
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        _set_run(p.add_run(), line, 60, TEXT, bold=True)
    # クリニック名（右下）
    cf = _textbox(slide, Inches(6), Inches(6.6), Inches(7), Inches(0.7))
    pc = cf.paragraphs[0]
    pc.alignment = PP_ALIGN.RIGHT
    _set_run(pc.add_run(), spec["clinic"], 20, TEXT, bold=True)


def _add_agenda_slide(prs, spec):
    slide = _blank_slide(prs)
    hf = _textbox(slide, Inches(1), Inches(1), Inches(6), Inches(2))
    _set_run(hf.paragraphs[0].add_run(), spec["title"], 44, TEXT, bold=True)
    bf = _textbox(slide, Inches(6.5), Inches(2.2), Inches(6), Inches(4))
    for i, item in enumerate(spec["items"]):
        p = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
        p.space_after = Pt(18)
        _set_run(p.add_run(), f"●  {item}", 28, TEXT)


def _add_text_slide(prs, spec):
    slide = _blank_slide(prs)
    hf = _textbox(slide, Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.2))
    _set_run(hf.paragraphs[0].add_run(), spec["title"], 38, TEXT, bold=True)

    bf = _textbox(slide, Inches(0.9), Inches(1.7), Inches(11.5), Inches(5.4))
    first = True
    for para in spec.get("body", []):
        p = bf.paragraphs[0] if first else bf.add_paragraph()
        first = False
        p.space_after = Pt(12)
        _set_run(p.add_run(), para, 20, TEXT)
    for heading, text in spec.get("blocks", []):
        ph = bf.paragraphs[0] if first else bf.add_paragraph()
        first = False
        ph.space_before = Pt(10)
        ph.space_after = Pt(4)
        _set_run(ph.add_run(), f"【{heading}】", 22, ACCENT, bold=True)
        pt = bf.add_paragraph()
        pt.space_after = Pt(12)
        _set_run(pt.add_run(), text, 20, TEXT)


def _add_photo_slide(prs, spec, images: dict[str, Image.Image], tmpdir: Path):
    slide = _blank_slide(prs)
    img = images.get(spec["key"])
    if img is None:
        # 画像が無い場合は見出しだけ置く
        hf = _textbox(slide, Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.2))
        _set_run(hf.paragraphs[0].add_run(),
                 f"{spec['title']}（{spec.get('subtitle', '')}）", 38, TEXT, bold=True)
        return
    # 生成画像（タイトル焼き込み済み・16:9）を全面に差し込む
    png = tmpdir / f"{spec['key']}.png"
    img.save(png)
    slide.shapes.add_picture(str(png), 0, 0, width=prs.slide_width, height=prs.slide_height)


def _add_closing_slide(prs, spec):
    slide = _blank_slide(prs)
    tf = _textbox(slide, Inches(1), Inches(3), Inches(11.3), Inches(1.5),
                  anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _set_run(p.add_run(), spec["title"], 54, TEXT, bold=True)


def build_pptx(images: dict[str, Image.Image], out_path: str | Path) -> Path:
    """DECK 構成に従って PPTX を生成する。

    images: スライドの key（composite_initial など）→ 生成済み PIL 画像。
    """
    out_path = Path(out_path)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    tmpdir = Path(tempfile.mkdtemp(prefix="hoteit_pptx_"))
    try:
        for spec in DECK:
            t = spec["type"]
            if t == "title":
                _add_title_slide(prs, spec)
            elif t == "agenda":
                _add_agenda_slide(prs, spec)
            elif t == "text":
                _add_text_slide(prs, spec)
            elif t == "photo":
                _add_photo_slide(prs, spec, images, tmpdir)
            elif t == "closing":
                _add_closing_slide(prs, spec)
        prs.save(out_path)
    finally:
        shutil.rmtree(tmpdir, ignore_errors=True)
    return out_path


def build_pdf(images: dict[str, Image.Image], out_path: str | Path) -> Path:
    """DECK を全スライド画像として描画し、1つの PDF にまとめる。

    LibreOffice に依存しないため、どの環境でも確実に PDF を生成できる。
    """
    out_path = Path(out_path)
    pages = render_deck(DECK, images)
    pages = [p.convert("RGB") if p.mode != "RGB" else p for p in pages]
    pages[0].save(
        out_path, format="PDF", save_all=True,
        append_images=pages[1:], resolution=150.0,
    )
    return out_path


def export_pdf(pptx_path: str | Path, out_dir: str | Path | None = None) -> Path | None:
    """LibreOffice を使って PPTX を PDF へ変換する。

    LibreOffice (soffice) が見つからない場合は None を返す。
    """
    pptx_path = Path(pptx_path)
    out_dir = Path(out_dir) if out_dir else pptx_path.parent

    soffice = shutil.which("libreoffice") or shutil.which("soffice")
    if not soffice:
        return None

    # 競合を避けるため一時プロファイルを使う
    with tempfile.TemporaryDirectory(prefix="lo_profile_") as profile:
        subprocess.run(
            [
                soffice, "--headless", "--convert-to", "pdf",
                "--outdir", str(out_dir),
                f"-env:UserInstallation=file://{profile}",
                str(pptx_path),
            ],
            check=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            timeout=180,
        )
    pdf_path = out_dir / (pptx_path.stem + ".pdf")
    return pdf_path if pdf_path.exists() else None
